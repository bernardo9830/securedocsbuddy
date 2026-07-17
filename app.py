# app.py - DocBuddy: web + agente con ricerca, riassunto Word, slide PPTX, citazione fonti e memoria.
import os
import uuid

from fastapi import FastAPI, HTTPException, Depends
from fastapi.responses import HTMLResponse, FileResponse
from pydantic import BaseModel

from dotenv import load_dotenv
from langchain_huggingface import HuggingFaceEndpointEmbeddings, HuggingFaceEndpoint, ChatHuggingFace
from langchain_community.vectorstores import FAISS
from langchain_core.tools import tool
from langchain.agents import create_agent
from langgraph.checkpoint.memory import InMemorySaver  # memoria della conversazione

from docx import Document as DocxDocument  # python-docx: crea i file Word
from pptx import Presentation              # python-pptx: crea le slide

load_dotenv()  # carica il token HF dal .env

CARTELLA_DOWNLOAD = "downloads"
os.makedirs(CARTELLA_DOWNLOAD, exist_ok=True)

# --- Componenti dell'agente: costruiti UNA sola volta, all'avvio ---
# Embeddings via HuggingFace Inference API (niente torch: leggero e cloud-ready).
embeddings = HuggingFaceEndpointEmbeddings(
    model="sentence-transformers/all-MiniLM-L6-v2",
    task="feature-extraction",
    huggingfacehub_api_token=os.getenv("HUGGINGFACEHUB_API_TOKEN"),
)
# Carico l'indice FAISS (leggero, niente chromadb/onnxruntime).
db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
retriever = db.as_retriever(search_kwargs={"k": 3})

motore = HuggingFaceEndpoint(repo_id="Qwen/Qwen2.5-7B-Instruct", task="text-generation")
llm = ChatHuggingFace(llm=motore)

# "Cassetto condiviso": segnala all'endpoint l'ultimo file creato (url + nome).
ultimo_file = {"url": None, "nome": None}

# Identificativo della conversazione (per la memoria). Cambiando questo si "svuota".
stato = {"thread": "conversazione-1"}


def _riassunto_da_documenti(argomento: str) -> str:
    """Funzione di supporto: recupera i documenti e produce un riassunto fedele."""
    docs = db.similarity_search(argomento, k=5)
    testo = "\n\n".join(d.page_content for d in docs)
    if not testo.strip():
        return ""
    prompt = (
        "Sei un assistente che riassume documenti. Leggi il TESTO e scrivi un riassunto "
        "FEDELE in italiano, basato SOLO sul testo (non inventare). "
        "Una frase di sintesi seguita da 4-6 punti elenco con le informazioni principali.\n\n"
        f"TESTO:\n{testo}\n\nRIASSUNTO:"
    )
    return llm.invoke(prompt).content


@tool
def cerca_documenti(domanda: str) -> str:
    """Cerca nei documenti aziendali le informazioni utili a rispondere a una domanda."""
    docs = retriever.invoke(domanda)
    return "\n\n".join(d.page_content for d in docs)


@tool
def crea_riassunto_word(argomento: str) -> str:
    """Crea un riassunto FEDELE dei documenti su un argomento e lo salva in un file Word scaricabile.
    Usa questo strumento quando l'utente chiede un riassunto, una sintesi o un documento Word."""
    riassunto = _riassunto_da_documenti(argomento)
    if not riassunto:
        return "Non ho trovato documenti pertinenti, quindi non ho creato il file."
    nome_file = f"riassunto_{uuid.uuid4().hex[:8]}.docx"
    percorso = os.path.join(CARTELLA_DOWNLOAD, nome_file)
    doc = DocxDocument()
    doc.add_heading(f"Riassunto: {argomento}", level=1)
    doc.add_paragraph(riassunto)
    doc.save(percorso)
    ultimo_file["url"] = f"/download/{nome_file}"
    ultimo_file["nome"] = nome_file
    return f"Ho creato un riassunto Word su '{argomento}', scaricabile."


@tool
def crea_slide_powerpoint(argomento: str) -> str:
    """Crea una presentazione PowerPoint di DUE slide che riassume i documenti su un argomento.
    Usa questo strumento quando l'utente chiede delle slide, una presentazione o un file PowerPoint."""
    riassunto = _riassunto_da_documenti(argomento)
    if not riassunto:
        return "Non ho trovato documenti pertinenti, quindi non ho creato le slide."

    prs = Presentation()
    # Slide 1: titolo
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = argomento
    s1.placeholders[1].text = "Riepilogo generato da DocBuddy"
    # Slide 2: contenuto (punti chiave)
    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "Punti chiave"
    s2.placeholders[1].text = riassunto

    nome_file = f"slide_{uuid.uuid4().hex[:8]}.pptx"
    percorso = os.path.join(CARTELLA_DOWNLOAD, nome_file)
    prs.save(percorso)
    ultimo_file["url"] = f"/download/{nome_file}"
    ultimo_file["nome"] = nome_file
    return f"Ho creato una presentazione di 2 slide su '{argomento}', scaricabile."


@tool
def cita_fonti(domanda: str) -> str:
    """Restituisce le fonti (i documenti) da cui proviene l'informazione per una domanda.
    Usa questo strumento quando l'utente chiede le fonti, le citazioni o 'da dove viene' una risposta."""
    docs = retriever.invoke(domanda)
    if not docs:
        return "Nessuna fonte trovata per questa domanda."
    righe = []
    for i, d in enumerate(docs, start=1):
        idx = d.metadata.get("id", "?")
        tema = d.metadata.get("domanda", "sconosciuto")
        estratto = d.page_content[:160].replace("\n", " ").strip()
        righe.append(f"[{i}] Documento #{idx} (tema: {tema})\n    \"{estratto}...\"")
    return "Fonti trovate:\n" + "\n".join(righe)


# La memoria: conserva i messaggi di ogni conversazione (in RAM).
memoria = InMemorySaver()

# L'agente ora ha QUATTRO strumenti e una MEMORIA.
agente = create_agent(
    model=llm,
    tools=[cerca_documenti, crea_riassunto_word, crea_slide_powerpoint, cita_fonti],
    checkpointer=memoria,
    system_prompt=(
        "Sei DocBuddy, un assistente sui documenti. "
        "Ricorda i messaggi precedenti e usali per le domande di follow-up. "
        "Usa cerca_documenti per rispondere a domande sui documenti. "
        "Usa crea_riassunto_word quando l'utente chiede un riassunto o un documento Word. "
        "Usa crea_slide_powerpoint quando l'utente chiede delle slide o una presentazione. "
        "Usa cita_fonti quando l'utente chiede le fonti o da dove viene un'informazione. "
        "Per i saluti rispondi direttamente senza usare strumenti."
    ),
)

app = FastAPI()

# Router dell'autenticazione (registrazione, login, /me)
from auth_router import router as auth_router
app.include_router(auth_router)

# Dependency e modello per proteggere gli endpoint e isolare per utente
from auth import get_current_user
from models import User

# Crea le tabelle del database all'avvio, se non esistono (idempotente): serve in cloud.
from database import Base, engine
Base.metadata.create_all(bind=engine)


class Richiesta(BaseModel):
    domanda: str


@app.get("/ping")
def ping():
    return {"messaggio": "pong"}


@app.post("/chiedi")
def chiedi(richiesta: Richiesta, utente: User = Depends(get_current_user)):
    # Depends(get_current_user): senza token valido FastAPI risponde 401 da solo.
    ultimo_file["url"] = None
    ultimo_file["nome"] = None
    # thread_id derivato dall'utente: ogni utente ha la SUA memoria isolata.
    config = {"configurable": {"thread_id": f"user-{utente.id}"}}
    risultato = agente.invoke(
        {"messages": [{"role": "user", "content": richiesta.domanda}]},
        config=config,
    )
    risposta = risultato["messages"][-1].content

    # Rete di sicurezza: se l'utente ha chiesto un file ma l'agente (modellino)
    # non ha chiamato il tool, lo invoco io in base alle parole chiave.
    testo = richiesta.domanda.lower()
    if ultimo_file["url"] is None:
        if any(k in testo for k in ["slide", "presentazione", "power point", "powerpoint", "ppt"]):
            risposta = crea_slide_powerpoint.invoke({"argomento": richiesta.domanda})
        elif any(k in testo for k in ["riassunt", "riassum", "sintesi", "word", "docx"]):
            risposta = crea_riassunto_word.invoke({"argomento": richiesta.domanda})

    return {"risposta": risposta, "download": ultimo_file["url"], "nome_file": ultimo_file["nome"]}


@app.post("/svuota-memoria")
def svuota_memoria(utente: User = Depends(get_current_user)):
    # Cancella SOLO la conversazione di questo utente.
    try:
        memoria.delete_thread(f"user-{utente.id}")
    except Exception:
        pass
    return {"ok": True}


@app.get("/download/{nome}")
def download(nome: str):
    nome = os.path.basename(nome)
    percorso = os.path.join(CARTELLA_DOWNLOAD, nome)
    if not os.path.exists(percorso):
        raise HTTPException(status_code=404, detail="File non trovato")
    return FileResponse(percorso, filename=nome)


PAGINA = """
<!DOCTYPE html>
<html lang="it">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>SecureDocs</title>
<style>
  * { box-sizing: border-box; }
  body {
    margin: 0;
    font-family: -apple-system, "Segoe UI", Roboto, Helvetica, Arial, sans-serif;
    background: linear-gradient(135deg, #667eea 0%, #764ba2 45%, #6a8dff 100%);
    background-size: 200% 200%; animation: sfondo 16s ease infinite;
    min-height: 100vh; display: flex; align-items: center; justify-content: center; padding: 20px;
  }
  @keyframes sfondo { 0%{background-position:0% 50%} 50%{background-position:100% 50%} 100%{background-position:0% 50%} }
  .card {
    position: relative; width: 100%; max-width: 580px; height: 82vh; max-height: 740px;
    background: #fff; border-radius: 20px; box-shadow: 0 24px 60px rgba(0,0,0,.3);
    display: flex; flex-direction: column; overflow: hidden;
  }
  header {
    background: linear-gradient(90deg, #4f46e5, #7c3aed); color: #fff; padding: 14px 18px;
    display: flex; align-items: center; gap: 10px;
  }
  .logo { width: 44px; height: 44px; background: #ffffff22; border-radius: 50%;
    display: flex; align-items: center; justify-content: center; font-size: 24px; }
  header h1 { margin: 0; font-size: 18px; }
  header p { margin: 2px 0 0; font-size: 12px; opacity: .85; }
  .header-btns { margin-left: auto; display: flex; gap: 6px; align-items: center; }
  .header-email { font-size: 12px; opacity: .9; margin-right: 4px; }
  .hbtn { background: #ffffff22; color: #fff; border: none; padding: 8px 12px;
    border-radius: 20px; cursor: pointer; font-size: 13px; }
  .hbtn:hover { background: #ffffff3a; }
  /* LOGIN */
  .login-body { flex: 1; display: flex; flex-direction: column; justify-content: center;
    gap: 12px; padding: 32px; }
  .login-body h2 { margin: 0 0 6px; color: #1f2937; }
  .login-body input { padding: 12px 14px; border: 1px solid #d1d5db; border-radius: 10px; font-size: 15px; outline: none; }
  .login-body input:focus { border-color: #7c3aed; box-shadow: 0 0 0 3px #7c3aed22; }
  .btn-primario { background: linear-gradient(135deg, #4f46e5, #7c3aed); color: #fff; border: none;
    padding: 12px; border-radius: 10px; font-size: 15px; cursor: pointer; }
  .btn-secondario { background: #fff; color: #4f46e5; border: 1px solid #c7d2fe;
    padding: 12px; border-radius: 10px; font-size: 15px; cursor: pointer; }
  .msg-err { color: #dc2626; font-size: 13px; min-height: 18px; }
  /* CHAT */
  #chat { flex: 1; overflow-y: auto; padding: 20px; background: #f4f5fb;
    display: flex; flex-direction: column; gap: 12px; }
  .riga { display: flex; gap: 8px; align-items: flex-end; max-width: 100%; }
  .riga-utente { justify-content: flex-end; }
  .riga-bot { justify-content: flex-start; }
  .avatar { width: 34px; height: 34px; border-radius: 50%; background: #ede9fe;
    display: flex; align-items: center; justify-content: center; font-size: 18px; flex-shrink: 0; }
  .msg { max-width: 80%; padding: 10px 14px; border-radius: 16px; line-height: 1.45;
    font-size: 15px; white-space: pre-wrap; word-wrap: break-word; box-shadow: 0 1px 2px rgba(0,0,0,.06); }
  .utente { background: linear-gradient(135deg, #4f46e5, #7c3aed); color: #fff; border-bottom-right-radius: 4px; }
  .bot { background: #fff; color: #1f2937; border: 1px solid #e5e7eb; border-bottom-left-radius: 4px; }
  .download-link { display: inline-block; margin-top: 8px; padding: 8px 12px; background: #10b981;
    color: #fff; border-radius: 8px; text-decoration: none; font-size: 14px; }
  .typing { display: flex; gap: 4px; align-items: center; }
  .typing span { width: 7px; height: 7px; border-radius: 50%; background: #9ca3af;
    display: inline-block; animation: blink 1.2s infinite both; }
  .typing span:nth-child(2) { animation-delay: .2s; }
  .typing span:nth-child(3) { animation-delay: .4s; }
  @keyframes blink { 0%,80%,100%{opacity:.3} 40%{opacity:1} }
  .barra { display: flex; gap: 8px; padding: 14px; background: #fff; border-top: 1px solid #e5e7eb; }
  #domanda { flex: 1; padding: 12px 16px; border: 1px solid #d1d5db; border-radius: 24px; font-size: 15px; outline: none; }
  #domanda:focus { border-color: #7c3aed; box-shadow: 0 0 0 3px #7c3aed22; }
  #invia { border: none; background: linear-gradient(135deg, #4f46e5, #7c3aed); color: #fff;
    width: 46px; height: 46px; border-radius: 50%; cursor: pointer;
    display: flex; align-items: center; justify-content: center; }
  #invia svg { width: 20px; height: 20px; }
  .pannello { position: absolute; inset: 0; background: #fff; display: flex; flex-direction: column;
    transform: translateX(100%); transition: transform .25s ease; z-index: 5; }
  .pannello.aperto { transform: translateX(0); }
  .pannello-head { background: linear-gradient(90deg, #4f46e5, #7c3aed); color: #fff; padding: 14px 18px;
    display: flex; align-items: center; justify-content: space-between; }
  .pannello-head button { background: #ffffff22; color: #fff; border: none; width: 30px; height: 30px;
    border-radius: 50%; cursor: pointer; }
  #storico-lista { flex: 1; overflow-y: auto; padding: 16px; }
  .voce { border-bottom: 1px solid #eee; padding: 10px 0; }
  .voce .q { font-weight: 600; color: #4f46e5; }
  .voce .a { color: #374151; margin-top: 4px; font-size: 14px; }
  .vuoto { color: #9ca3af; text-align: center; margin-top: 40px; }
</style>
</head>
<body>

  <!-- ===== SCHERMATA DI LOGIN ===== -->
  <div id="login-view" class="card">
    <header>
      <div class="logo">🤖</div>
      <div><h1>SecureDocs</h1><p>Accedi per usare l'assistente</p></div>
    </header>
    <div class="login-body">
      <h2>Accedi</h2>
      <input id="login-email" type="email" placeholder="Email">
      <input id="login-password" type="password" placeholder="Password">
      <div id="login-msg" class="msg-err"></div>
      <button class="btn-primario" onclick="accedi()">Accedi</button>
      <button class="btn-secondario" onclick="registra()">Registrati</button>
    </div>
  </div>

  <!-- ===== SCHERMATA CHAT (nascosta finche' non fai login) ===== -->
  <div id="chat-view" class="card" style="display:none">
    <header>
      <div class="logo">🤖</div>
      <div><h1>SecureDocs</h1><p>Assistente sui documenti</p></div>
      <div class="header-btns">
        <span id="header-email" class="header-email"></span>
        <button class="hbtn" onclick="toggleMemoria()">Memoria</button>
        <button class="hbtn" onclick="svuotaMemoria()">Svuota</button>
        <button class="hbtn" onclick="logout()">Esci</button>
      </div>
    </header>
    <div id="chat"></div>
    <div class="barra">
      <input id="domanda" placeholder="Scrivi una domanda...">
      <button id="invia" onclick="invia()" title="Invia">
        <svg viewBox="0 0 24 24" fill="none" stroke="white" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><line x1="22" y1="2" x2="11" y2="13"></line><polygon points="22 2 15 22 11 13 2 9 22 2"></polygon></svg>
      </button>
    </div>
    <div id="pannello" class="pannello">
      <div class="pannello-head"><strong>Risposte precedenti</strong><button onclick="toggleMemoria()">X</button></div>
      <div id="storico-lista"></div>
    </div>
  </div>

<script>
  let token = localStorage.getItem("sd_token") || null;
  let storico = JSON.parse(localStorage.getItem("sd_storico") || "[]");

  const loginView = document.getElementById("login-view");
  const chatView = document.getElementById("chat-view");
  const chat = document.getElementById("chat");

  function mostraChat() {
    loginView.style.display = "none";
    chatView.style.display = "flex";
    if (chat.childElementCount === 0) benvenuto();
    caricaEmail();
  }
  function mostraLogin() {
    chatView.style.display = "none";
    loginView.style.display = "flex";
  }

  // --- AUTENTICAZIONE ---
  async function accedi() {
    const email = document.getElementById("login-email").value.trim();
    const password = document.getElementById("login-password").value;
    const msg = document.getElementById("login-msg");
    msg.textContent = "";
    try {
      const resp = await fetch("/login", {
        method: "POST",
        headers: { "Content-Type": "application/x-www-form-urlencoded" },
        body: new URLSearchParams({ username: email, password: password })
      });
      if (!resp.ok) { msg.textContent = "Email o password errati."; return; }
      const dati = await resp.json();
      token = dati.access_token;
      localStorage.setItem("sd_token", token);
      mostraChat();
    } catch (e) { msg.textContent = "Errore di rete: " + e; }
  }

  async function registra() {
    const email = document.getElementById("login-email").value.trim();
    const password = document.getElementById("login-password").value;
    const msg = document.getElementById("login-msg");
    msg.textContent = "";
    if (!email || !password) { msg.textContent = "Inserisci email e password."; return; }
    try {
      const resp = await fetch("/registrazione", {
        method: "POST",
        headers: { "Content-Type": "application/json" },
        body: JSON.stringify({ email: email, password: password })
      });
      if (resp.status === 201) { await accedi(); return; }
      if (resp.status === 400) { msg.textContent = "Email gia' registrata: prova ad accedere."; return; }
      msg.textContent = "Registrazione non riuscita.";
    } catch (e) { msg.textContent = "Errore di rete: " + e; }
  }

  async function caricaEmail() {
    try {
      const resp = await fetch("/me", { headers: { "Authorization": "Bearer " + token } });
      if (resp.ok) { const u = await resp.json(); document.getElementById("header-email").textContent = u.email; }
    } catch (e) {}
  }

  function logout() {
    token = null;
    localStorage.removeItem("sd_token");
    mostraLogin();
  }

  // --- CHAT ---
  function rigaUtente(testo) {
    const riga = document.createElement("div"); riga.className = "riga riga-utente";
    const b = document.createElement("div"); b.className = "msg utente"; b.textContent = testo;
    riga.appendChild(b); chat.appendChild(riga); chat.scrollTop = chat.scrollHeight;
  }
  function rigaBot() {
    const riga = document.createElement("div"); riga.className = "riga riga-bot";
    const av = document.createElement("div"); av.className = "avatar"; av.textContent = "🤖";
    const b = document.createElement("div"); b.className = "msg bot";
    riga.appendChild(av); riga.appendChild(b); chat.appendChild(riga); chat.scrollTop = chat.scrollHeight;
    return b;
  }
  function benvenuto() {
    const b = rigaBot();
    b.textContent = "Ciao! Sono SecureDocs. Chiedimi qualcosa, un riassunto, delle slide o le fonti.";
  }

  async function invia() {
    const input = document.getElementById("domanda");
    const domanda = input.value.trim();
    if (!domanda) return;
    rigaUtente(domanda); input.value = "";
    const bolla = rigaBot(); bolla.className = "msg bot typing";
    bolla.innerHTML = "<span></span><span></span><span></span>";
    try {
      const resp = await fetch("/chiedi", {
        method: "POST",
        headers: { "Content-Type": "application/json", "Authorization": "Bearer " + token },
        body: JSON.stringify({ domanda: domanda })
      });
      if (resp.status === 401) { logout(); return; }
      const dati = await resp.json();
      bolla.className = "msg bot"; bolla.textContent = dati.risposta;
      if (dati.download) {
        const a = document.createElement("a");
        a.href = dati.download; a.className = "download-link";
        a.textContent = "Scarica " + (dati.nome_file || "il file");
        a.setAttribute("download", "");
        bolla.appendChild(document.createElement("br")); bolla.appendChild(a);
      }
      storico.push({ domanda: domanda, risposta: dati.risposta });
      localStorage.setItem("sd_storico", JSON.stringify(storico));
    } catch (e) {
      bolla.className = "msg bot"; bolla.textContent = "Errore: " + e;
    } finally { input.focus(); chat.scrollTop = chat.scrollHeight; }
  }

  async function svuotaMemoria() {
    if (!confirm("Vuoi svuotare la memoria della conversazione?")) return;
    try {
      await fetch("/svuota-memoria", { method: "POST", headers: { "Authorization": "Bearer " + token } });
    } catch (e) {}
    storico = []; localStorage.removeItem("sd_storico");
    chat.innerHTML = ""; benvenuto();
    document.getElementById("pannello").classList.remove("aperto");
  }

  function toggleMemoria() {
    const pannello = document.getElementById("pannello");
    const lista = document.getElementById("storico-lista");
    if (storico.length === 0) { lista.innerHTML = '<div class="vuoto">Nessuna risposta ancora.</div>'; }
    else {
      lista.innerHTML = "";
      storico.forEach(function (v) {
        const div = document.createElement("div"); div.className = "voce";
        const q = document.createElement("div"); q.className = "q"; q.textContent = "D: " + v.domanda;
        const a = document.createElement("div"); a.className = "a"; a.textContent = v.risposta;
        div.appendChild(q); div.appendChild(a); lista.appendChild(div);
      });
    }
    pannello.classList.toggle("aperto");
  }

  document.getElementById("domanda").addEventListener("keydown", function (e) { if (e.key === "Enter") invia(); });
  document.getElementById("login-password").addEventListener("keydown", function (e) { if (e.key === "Enter") accedi(); });

  // All'avvio: se ho gia' un token mostro la chat, altrimenti il login.
  if (token) { mostraChat(); } else { mostraLogin(); }
</script>
</body>
</html>
"""


@app.get("/", response_class=HTMLResponse)
def home():
    return PAGINA
